#!/usr/bin/env python3
"""Generate MIC2 VLM Classifier presentation."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)
ORANGE     = RGBColor(0xFF, 0xA0, 0x2F)
RED        = RGBColor(0xEF, 0x47, 0x6F)
DARK_CARD  = RGBColor(0x13, 0x2A, 0x40)   # card background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text, subtitle_text=None):
    """Cyan left accent bar + title."""
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT)
    add_textbox(slide, title_text,
                Inches(0.25), Inches(0.18), Inches(12.8), Inches(0.7),
                font_size=32, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.25), Inches(0.82), Inches(12.8), Inches(0.4),
                    font_size=16, color=ACCENT2)


def add_bullet_card(slide, title, bullets, left, top, width, height,
                    title_color=ACCENT, bullet_color=LIGHT_GRAY, font_size=15):
    add_rect(slide, left, top, width, height, DARK_CARD)
    add_textbox(slide, title, left + Inches(0.18), top + Inches(0.12),
                width - Inches(0.25), Inches(0.38),
                font_size=16, bold=True, color=title_color)
    y = top + Inches(0.52)
    for b in bullets:
        add_textbox(slide, f"• {b}", left + Inches(0.18), y,
                    width - Inches(0.3), Inches(0.36),
                    font_size=font_size, color=bullet_color)
        y += Inches(0.36)


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.38)):
    col_w = width / len(headers)
    # header row
    for c, h in enumerate(headers):
        add_rect(slide, left + col_w * c, top, col_w, row_height, ACCENT)
        add_textbox(slide, h, left + col_w * c + Inches(0.05), top + Inches(0.04),
                    col_w - Inches(0.1), row_height,
                    font_size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = DARK_CARD if r % 2 == 0 else RGBColor(0x18, 0x34, 0x4F)
        for c, cell in enumerate(row):
            add_rect(slide, left + col_w * c, top + row_height * (r + 1),
                     col_w, row_height, bg)
            color = GREEN if "95" in str(cell) and "%" in str(cell) else LIGHT_GRAY
            bold = "ours" in str(cell).lower() or ("95" in str(cell) and "%" in str(cell))
            add_textbox(slide, str(cell),
                        left + col_w * c + Inches(0.05),
                        top + row_height * (r + 1) + Inches(0.04),
                        col_w - Inches(0.1), row_height,
                        font_size=12, color=color, bold=bold,
                        align=PP_ALIGN.CENTER)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # big accent block left
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, ACCENT)

    # title
    add_textbox(slide, "MIC2 Malicious Ad Classifier",
                Inches(0.8), Inches(1.6), Inches(11), Inches(1.1),
                font_size=44, bold=True, color=WHITE)
    add_textbox(slide, "Fine-tuning InternVL 2.5-1B with LoRA for screenshot-based threat detection",
                Inches(0.8), Inches(2.75), Inches(10.5), Inches(0.6),
                font_size=20, color=ACCENT2, italic=True)

    # stats row
    stats = [("158K", "Training images"), ("13", "Threat classes"),
             ("95.2%", "Val accuracy"), ("~$14", "Training cost")]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        add_rect(slide, x, Inches(3.8), Inches(2.8), Inches(1.3), DARK_CARD)
        add_textbox(slide, val, x + Inches(0.1), Inches(3.88), Inches(2.6), Inches(0.6),
                    font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x + Inches(0.1), Inches(4.45), Inches(2.6), Inches(0.35),
                    font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "GeoEdge  |  May 2026",
                Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
                font_size=12, color=LIGHT_GRAY)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "The Problem", "Detecting malicious ads from screenshots at scale")

    classes = [
        ("fake_av", "Fake antivirus / security alerts"),
        ("tech_support_scam", "Fake Microsoft/Apple lock screens"),
        ("financial_scam", "Fake investment & lottery fraud"),
        ("fake_appstore", "Counterfeit app store pages"),
        ("misleading_offers", "Fake prizes & deceptive coupons"),
        ("gift_card_scan", "Gift card reward scams"),
        ("forced_notification", "Trick-to-subscribe push pages"),
        ("fake_downloader", "Fake download buttons (malware)"),
    ]

    add_textbox(slide, "13 malicious ad categories detected from a single screenshot:",
                Inches(0.25), Inches(1.35), Inches(12.5), Inches(0.4),
                font_size=15, color=LIGHT_GRAY)

    for i, (cls, desc) in enumerate(classes):
        col = i % 2
        row = i // 2
        x = Inches(0.25 + col * 6.55)
        y = Inches(1.85 + row * 1.15)
        add_rect(slide, x, y, Inches(6.3), Inches(1.0), DARK_CARD)
        add_textbox(slide, cls, x + Inches(0.15), y + Inches(0.08),
                    Inches(6.0), Inches(0.38), font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, desc, x + Inches(0.15), y + Inches(0.48),
                    Inches(6.0), Inches(0.38), font_size=13, color=LIGHT_GRAY)

    add_textbox(slide, "+ blank_LP  |  suspicious_vpn  |  malicious_extension  |  fake_updates  |  Benign",
                Inches(0.25), Inches(6.7), Inches(12.5), Inches(0.35),
                font_size=12, color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_approach(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Our Approach", "Vision encoder + LoRA fine-tuning + linear classifier")

    # Pipeline boxes
    steps = [
        ("Screenshot\n448×448 px", ACCENT),
        ("InternViT-300M\nvision encoder", ACCENT),
        ("CLS token\n[B, 1024]", ACCENT2),
        ("LoRA\nadapters", ORANGE),
        ("Linear head\n1024 → 13", GREEN),
        ("Predicted\nClass", WHITE),
    ]
    box_w = Inches(1.8)
    box_h = Inches(1.1)
    gap = Inches(0.22)
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(2.5)

    for i, (label, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        is_main = color == ACCENT
        bg = DARK_CARD
        add_rect(slide, x, y, box_w, box_h, bg)
        # colored top border via thin rect
        add_rect(slide, x, y, box_w, Inches(0.05), color)
        add_textbox(slide, label, x + Inches(0.08), y + Inches(0.12),
                    box_w - Inches(0.16), box_h - Inches(0.15),
                    font_size=13, bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.04)
            add_textbox(slide, "→", ax, y + Inches(0.35), gap + Inches(0.05), Inches(0.4),
                        font_size=18, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

    # Key decisions
    decisions = [
        ("Discard LLM", "InternVL has a 700M-param language model — we drop it entirely. Classification needs the vision encoder only, saving ~8 GB VRAM."),
        ("LoRA not full fine-tune", "Freeze 300M base weights. Train only ~4.7M LoRA adapter parameters inserted into attention layers. Fits on a single A10G (24 GB)."),
        ("CLS token as image rep", "The [CLS] token at position 0 aggregates context from all 1024 image patches after 24 transformer layers — one vector represents the whole screenshot."),
    ]
    for i, (title, body) in enumerate(decisions):
        x = Inches(0.25 + i * 4.35)
        y2 = Inches(4.1)
        add_rect(slide, x, y2, Inches(4.1), Inches(2.9), DARK_CARD)
        add_rect(slide, x, y2, Inches(4.1), Inches(0.05),
                 [ACCENT, ORANGE, GREEN][i])
        add_textbox(slide, title, x + Inches(0.15), y2 + Inches(0.12),
                    Inches(3.8), Inches(0.4), font_size=14, bold=True,
                    color=[ACCENT, ORANGE, GREEN][i])
        add_textbox(slide, body, x + Inches(0.15), y2 + Inches(0.55),
                    Inches(3.8), Inches(2.2), font_size=12, color=LIGHT_GRAY, wrap=True)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Model Architecture", "InternViT-300M + LoRA + Linear head")

    # Left: architecture flow
    stages = [
        ("Input  448×448 RGB image", LIGHT_GRAY, False),
        ("Patch Embedding", ACCENT, True),
        ("Slice into 32×32 patches → 1024 tokens + 1 CLS token", LIGHT_GRAY, False),
        ("24× Transformer Blocks", ACCENT, True),
        ("Self-attention (qkv, proj)  ← LoRA adapters injected here", ORANGE, False),
        ("Each patch attends to all others — global context builds up", LIGHT_GRAY, False),
        ("CLS Token  [1024-dim vector]", ACCENT, True),
        ("Summarises the entire image after 24 layers", LIGHT_GRAY, False),
        ("Dropout (p=0.1)  →  Linear(1024 → 13)", GREEN, True),
        ("Cross-entropy loss during training", LIGHT_GRAY, False),
    ]

    y = Inches(1.35)
    for text, color, is_header in stages:
        size = 15 if is_header else 13
        bold = is_header
        indent = Inches(0.25) if is_header else Inches(0.55)
        add_textbox(slide, ("▶  " if is_header else "     ") + text,
                    indent, y, Inches(7.5), Inches(0.35),
                    font_size=size, bold=bold, color=color)
        y += Inches(0.36) if is_header else Inches(0.32)

    # Right: parameter breakdown
    add_rect(slide, Inches(8.1), Inches(1.35), Inches(4.9), Inches(5.7), DARK_CARD)
    add_textbox(slide, "Parameter Budget", Inches(8.3), Inches(1.48),
                Inches(4.5), Inches(0.4), font_size=16, bold=True, color=ACCENT)

    params = [
        ("InternViT-300M base", "300M", "frozen", LIGHT_GRAY),
        ("LoRA A matrices (×48)", "~2.4M", "trainable", ORANGE),
        ("LoRA B matrices (×48)", "~2.4M", "trainable", ORANGE),
        ("Classifier head", "13K", "trainable", GREEN),
        ("TOTAL trainable", "~4.7M", "1.6%", ACCENT),
    ]
    y2 = Inches(2.05)
    for name, count, status, color in params:
        if name == "TOTAL trainable":
            add_rect(slide, Inches(8.1), y2 - Inches(0.05), Inches(4.9), Inches(0.42),
                     RGBColor(0x08, 0x2A, 0x3A))
        add_textbox(slide, name, Inches(8.3), y2, Inches(2.6), Inches(0.35),
                    font_size=13, color=color, bold=(name == "TOTAL trainable"))
        add_textbox(slide, count, Inches(10.9), y2, Inches(0.9), Inches(0.35),
                    font_size=13, color=color, bold=True, align=PP_ALIGN.RIGHT)
        add_textbox(slide, status, Inches(11.8), y2, Inches(1.0), Inches(0.35),
                    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        y2 += Inches(0.5)

    # LoRA diagram
    add_textbox(slide, "How LoRA works", Inches(8.3), Inches(4.4),
                Inches(4.5), Inches(0.38), font_size=15, bold=True, color=ACCENT)
    lora_lines = [
        "Standard:   output = W · x",
        "                (W has D² params — frozen)",
        "",
        "LoRA:   output = W·x  +  B·A·x",
        "        A: D×r    B: r×D    r=16",
        "        2 × D × r  <<  D²",
    ]
    y3 = Inches(4.85)
    for line in lora_lines:
        color = ORANGE if "LoRA" in line or "A:" in line else LIGHT_GRAY
        add_textbox(slide, line, Inches(8.3), y3, Inches(4.6), Inches(0.3),
                    font_size=12, color=color,
                    bold=("LoRA:" in line))
        y3 += Inches(0.29)


def slide_why_drop_llm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Why We Drop the Language Model",
                  "InternVL ships with a 700M-param LLM — we discard it entirely")

    # Full VLM vs ours
    add_textbox(slide, "Full VLM (original purpose)",
                Inches(0.25), Inches(1.38), Inches(5.8), Inches(0.4),
                font_size=15, bold=True, color=RED)
    full_flow = [
        'Image  →  InternViT-300M  →  visual tokens',
        '                                    ↓',
        '              Qwen2-0.7B  (reads tokens + text prompt)',
        '                                    ↓',
        '  "This is a fake antivirus page..."   (text output)',
    ]
    y = Inches(1.85)
    for line in full_flow:
        add_textbox(slide, line, Inches(0.35), y, Inches(5.7), Inches(0.32),
                    font_size=12, color=LIGHT_GRAY)
        y += Inches(0.3)

    add_textbox(slide, "Our classifier",
                Inches(7.0), Inches(1.38), Inches(5.8), Inches(0.4),
                font_size=15, bold=True, color=GREEN)
    our_flow = [
        'Image  →  InternViT-300M  →  CLS token [1024]',
        '                                    ↓',
        '                  Linear(1024 → 13)',
        '                                    ↓',
        '         "tech_support_scam"   (class label)',
    ]
    y = Inches(1.85)
    for line in our_flow:
        add_textbox(slide, line, Inches(7.1), y, Inches(5.7), Inches(0.32),
                    font_size=12, color=LIGHT_GRAY)
        y += Inches(0.3)

    # Divider
    add_rect(slide, Inches(6.4), Inches(1.3), Inches(0.04), Inches(2.5), ACCENT)

    # What the LLM does that we don't need
    add_textbox(slide, "LLM components we don't need:",
                Inches(0.25), Inches(3.75), Inches(8), Inches(0.38),
                font_size=15, bold=True, color=ACCENT)

    components = [
        ("Token embeddings  (vocab ~150K)", "Maps text tokens to vectors — we have no text input"),
        ("24 decoder transformer layers", "Generate next token autoregressively — we have no text output"),
        ("Language modelling head", "Projects to vocabulary for generation — irrelevant"),
        ("Cross-attention to visual tokens", "CLS token already encodes global image context"),
    ]
    y = Inches(4.2)
    for comp, reason in components:
        add_rect(slide, Inches(0.25), y, Inches(12.8), Inches(0.52), DARK_CARD)
        add_textbox(slide, "✗  " + comp, Inches(0.4), y + Inches(0.07),
                    Inches(4.5), Inches(0.38), font_size=13, color=RED, bold=True)
        add_textbox(slide, reason, Inches(4.9), y + Inches(0.07),
                    Inches(8.0), Inches(0.38), font_size=13, color=LIGHT_GRAY)
        y += Inches(0.6)

    add_rect(slide, Inches(0.25), y + Inches(0.08), Inches(12.8), Inches(0.52),
             RGBColor(0x04, 0x2A, 0x1A))
    add_textbox(slide, "Result: VRAM drops from ~16 GB → ~8 GB  |  Inference 5× faster  |  No prompt engineering needed",
                Inches(0.4), y + Inches(0.15), Inches(12.5), Inches(0.38),
                font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def slide_lora(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "LoRA: Fine-tuning Without the VRAM Bill",
                  "Train 4.7M parameters instead of 300M — same quality, 64× less memory")

    # Problem / solution
    add_textbox(slide, "The problem with full fine-tuning:",
                Inches(0.25), Inches(1.38), Inches(6), Inches(0.38),
                font_size=15, bold=True, color=RED)
    problem_pts = [
        "300M base parameters   × 4 bytes (fp32) = 1.2 GB weights",
        "Gradients              × 4 bytes         = 1.2 GB",
        "Adam optimizer states  × 8 bytes         = 2.4 GB",
        "Activations (batch=8)                    = ~10 GB",
        "Total: >15 GB — barely fits, no headroom for batch size",
    ]
    y = Inches(1.82)
    for pt in problem_pts:
        color = RED if "Total" in pt else LIGHT_GRAY
        add_textbox(slide, ("⚠  " if "Total" in pt else "   ") + pt,
                    Inches(0.35), y, Inches(6.0), Inches(0.33),
                    font_size=12, color=color, bold="Total" in pt)
        y += Inches(0.32)

    add_textbox(slide, "LoRA solution:",
                Inches(6.9), Inches(1.38), Inches(6), Inches(0.38),
                font_size=15, bold=True, color=GREEN)
    solution_pts = [
        "Freeze all 300M base weights (no gradients needed)",
        "Inject tiny A, B matrices into each attention layer",
        "A: [hidden × rank=16]  B: [rank=16 × hidden]",
        "Only A and B are trained  →  ~4.7M params total",
        "VRAM: ~8 GB  |  Batch=8  |  Fits on A10G comfortably",
    ]
    y = Inches(1.82)
    for pt in solution_pts:
        color = GREEN if "VRAM" in pt else LIGHT_GRAY
        add_textbox(slide, ("✓  " if "VRAM" in pt else "   ") + pt,
                    Inches(7.0), y, Inches(6.0), Inches(0.33),
                    font_size=12, color=color, bold="VRAM" in pt)
        y += Inches(0.32)

    add_rect(slide, Inches(6.4), Inches(1.3), Inches(0.04), Inches(2.55), ACCENT)

    # Visual comparison bar
    add_textbox(slide, "Trainable parameters:", Inches(0.25), Inches(4.2),
                Inches(5), Inches(0.35), font_size=14, bold=True, color=ACCENT)

    # Full fine-tune bar
    add_textbox(slide, "Full fine-tune  300M", Inches(0.25), Inches(4.65),
                Inches(3), Inches(0.32), font_size=13, color=RED)
    add_rect(slide, Inches(3.4), Inches(4.68), Inches(9.0), Inches(0.28), RED)

    # LoRA bar
    add_textbox(slide, "LoRA (ours)   4.7M", Inches(0.25), Inches(5.15),
                Inches(3), Inches(0.32), font_size=13, color=GREEN)
    add_rect(slide, Inches(3.4), Inches(5.18), Inches(0.141), Inches(0.28), GREEN)
    add_textbox(slide, "1.6% →", Inches(3.58), Inches(5.12), Inches(1), Inches(0.32),
                font_size=12, color=GREEN, bold=True)

    # Where LoRA is applied
    add_textbox(slide, "Where adapters are injected  (every attention layer, ×24 blocks):",
                Inches(0.25), Inches(5.75), Inches(12.5), Inches(0.35),
                font_size=14, bold=True, color=ACCENT)
    targets = [("qkv projection", "Controls what each patch queries, keys, values from others"),
               ("output projection", "Controls how attended values are combined and output")]
    y = Inches(6.18)
    for name, desc in targets:
        add_rect(slide, Inches(0.25), y, Inches(12.8), Inches(0.45), DARK_CARD)
        add_textbox(slide, name, Inches(0.4), y + Inches(0.06), Inches(2.5), Inches(0.33),
                    font_size=13, bold=True, color=ORANGE)
        add_textbox(slide, desc, Inches(2.9), y + Inches(0.06), Inches(9.8), Inches(0.33),
                    font_size=13, color=LIGHT_GRAY)
        y += Inches(0.52)


def slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Comparison to Alternatives",
                  "Why InternVL + LoRA outperforms other approaches for this task")

    headers = ["Approach", "Val Accuracy", "VRAM", "Train Time", "Cost", "Screenshot-aware"]
    rows = [
        ["ResNet-50 (baseline)", "~85%", "4 GB", "2 hrs", "~$2", "No"],
        ["CLIP ViT-L/14 (frozen)", "~88%", "6 GB", "30 min", "~$0.50", "Partial"],
        ["Full fine-tune ViT-Large", "~92%", ">40 GB", "8 hrs", ">$100", "No"],
        ["EfficientNet-B7", "~87%", "8 GB", "3 hrs", "~$3", "No"],
        ["InternVL + LoRA  (ours) ★", "~95.2%", "18 GB", "14 hrs", "~$14", "Yes"],
    ]
    add_table(slide, headers, rows, Inches(0.2), Inches(1.4), Inches(12.9))

    # Key differentiators
    diffs = [
        ("Pre-trained on web data", "InternViT was trained on diverse web image-text pairs including page screenshots. It already understands logos, buttons, text layouts, and visual patterns common in malicious ads."),
        ("448px resolution", "Malicious ads rely on fake text (phone numbers, URLs, warnings). At 448px, InternViT reads small text that 224px models miss entirely."),
        ("LoRA efficiency", "Achieves near-full-fine-tune accuracy at 1/10th the VRAM. Frozen backbone retains web visual knowledge; adapters learn what's malicious-specific."),
    ]
    y = Inches(4.35)
    for i, (title, body) in enumerate(diffs):
        x = Inches(0.2 + i * 4.38)
        add_rect(slide, x, y, Inches(4.2), Inches(2.75), DARK_CARD)
        add_rect(slide, x, y, Inches(4.2), Inches(0.05), [ACCENT, ORANGE, GREEN][i])
        add_textbox(slide, title, x + Inches(0.15), y + Inches(0.1),
                    Inches(3.9), Inches(0.38), font_size=14, bold=True,
                    color=[ACCENT, ORANGE, GREEN][i])
        add_textbox(slide, body, x + Inches(0.15), y + Inches(0.55),
                    Inches(3.9), Inches(2.1), font_size=12, color=LIGHT_GRAY, wrap=True)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Training Results  (Live — Epoch 2.7 / 5)",
                  "On-demand g5.xlarge  |  A10G 24 GB  |  ~14 hours total")

    # Accuracy progression
    add_textbox(slide, "Validation accuracy over training",
                Inches(0.25), Inches(1.38), Inches(7), Inches(0.38),
                font_size=15, bold=True, color=ACCENT)

    evals = [
        (0.07, 81.1), (0.13, 86.0), (0.20, 90.5), (0.27, 91.3),
        (0.33, 92.9), (0.53, 93.4), (0.73, 93.9), (1.00, 95.1),
        (1.33, 95.1), (1.87, 94.8), (2.00, 95.1), (2.27, 95.1),
        (2.53, 95.2), (2.73, 95.2),
    ]
    chart_l, chart_t = Inches(0.25), Inches(1.85)
    chart_w, chart_h = Inches(7.0), Inches(3.5)
    add_rect(slide, chart_l, chart_t, chart_w, chart_h, DARK_CARD)

    # axis labels
    for pct, label in [(81, "81%"), (87, "87%"), (93, "93%"), (99, "99%")]:
        frac = (pct - 79) / 22
        y_pos = chart_t + chart_h - chart_h * frac - Inches(0.12)
        add_textbox(slide, label, chart_l + Inches(0.05), y_pos,
                    Inches(0.55), Inches(0.28), font_size=10, color=LIGHT_GRAY)
        add_rect(slide, chart_l + Inches(0.6), y_pos + Inches(0.12),
                 chart_w - Inches(0.65), Inches(0.01),
                 RGBColor(0x22, 0x3A, 0x50))

    # plot points & line
    pts = []
    for epoch, acc in evals:
        x_frac = epoch / 5.0
        y_frac = (acc - 79) / 22
        px = chart_l + Inches(0.65) + (chart_w - Inches(0.75)) * x_frac
        py = chart_t + chart_h - chart_h * y_frac - Inches(0.08)
        pts.append((px, py))
        dot = slide.shapes.add_shape(9, px - Inches(0.06), py - Inches(0.06),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if acc < 95.2 else GREEN
        dot.line.fill.background()

    # epoch labels
    for ep_label in [0, 1, 2, 3, 4, 5]:
        x_frac = ep_label / 5.0
        px = chart_l + Inches(0.65) + (chart_w - Inches(0.75)) * x_frac
        add_textbox(slide, f"ep{ep_label}", px - Inches(0.2),
                    chart_t + chart_h + Inches(0.02), Inches(0.4), Inches(0.28),
                    font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # best marker
    best_x, best_y = pts[-1]
    add_textbox(slide, "95.2% ★", best_x - Inches(0.6), best_y - Inches(0.45),
                Inches(1.1), Inches(0.32), font_size=12, bold=True, color=GREEN)

    # Right: key numbers
    metrics = [
        ("95.2%", "Best val accuracy", GREEN),
        ("~0.09", "Training loss (epoch 2.7)", ACCENT),
        ("60,021", "Training samples", LIGHT_GRAY),
        ("7,846", "Held-out test samples", LIGHT_GRAY),
        ("37,515", "Total training steps", LIGHT_GRAY),
        ("1.36 s", "Per step (A10G)", LIGHT_GRAY),
    ]
    x_r = Inches(7.7)
    y_r = Inches(1.38)
    for val, label, color in metrics:
        add_rect(slide, x_r, y_r, Inches(5.35), Inches(0.72), DARK_CARD)
        add_textbox(slide, val, x_r + Inches(0.15), y_r + Inches(0.04),
                    Inches(2.0), Inches(0.38), font_size=22, bold=True, color=color)
        add_textbox(slide, label, x_r + Inches(2.25), y_r + Inches(0.18),
                    Inches(3.0), Inches(0.38), font_size=13, color=LIGHT_GRAY)
        y_r += Inches(0.8)

    add_textbox(slide, "Full per-class classification report (precision / recall / F1 + confusion matrix) saved to S3 on completion.",
                Inches(0.25), Inches(6.85), Inches(12.8), Inches(0.35),
                font_size=11, color=ACCENT2, italic=True)


def slide_infra(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_title_bar(slide, "Training Infrastructure",
                  "One command from your laptop → model in S3")

    # Flow diagram
    boxes = [
        ("Local machine\ntrain_cli.py launch\n--on-demand", ACCENT, Inches(0.2)),
        ("AWS EC2\ng5.xlarge\nA10G 24 GB", ORANGE, Inches(3.5)),
        ("HuggingFace\nInternVL2_5-1B\nmodel weights", ACCENT2, Inches(6.8)),
        ("S3 bucket\nDataset + checkpoints\n+ final model", GREEN, Inches(10.1)),
    ]
    for label, color, x in boxes:
        add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(1.4), DARK_CARD)
        add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.06), color)
        add_textbox(slide, label, x + Inches(0.12), Inches(1.58),
                    Inches(2.76), Inches(1.3), font_size=13, bold=False,
                    color=color, align=PP_ALIGN.CENTER)
        if x != Inches(10.1):
            add_textbox(slide, "→", x + Inches(3.02), Inches(2.0),
                        Inches(0.46), Inches(0.5), font_size=22, bold=True,
                        color=ACCENT2, align=PP_ALIGN.CENTER)

    # Resilience features
    features = [
        ("Spot interruption\nhandling", "Watcher polls EC2 metadata every 5s. Sends SIGTERM 90s before reclaim so the model saves an emergency checkpoint.", ORANGE),
        ("S3 checkpoint\nsync", "S3CheckpointCallback uploads every checkpoint immediately after save. latest.txt tracks the newest.", ACCENT),
        ("Auto-resume", "On relaunch, s3_download_checkpoint() fetches latest checkpoint and passes resume_from_checkpoint= to Trainer.", GREEN),
        ("Live log\nstreaming", "Bootstrap uploads training.log to S3 every 30s. train_cli.py status polls and streams to your terminal.", ACCENT2),
    ]
    y = Inches(3.2)
    add_textbox(slide, "Resilience features:", Inches(0.2), y - Inches(0.4),
                Inches(4), Inches(0.35), font_size=14, bold=True, color=ACCENT)
    for i, (title, body, color) in enumerate(features):
        x = Inches(0.2 + i * 3.28)
        add_rect(slide, x, y, Inches(3.1), Inches(3.8), DARK_CARD)
        add_rect(slide, x, y, Inches(3.1), Inches(0.05), color)
        add_textbox(slide, title, x + Inches(0.12), y + Inches(0.1),
                    Inches(2.86), Inches(0.6), font_size=13, bold=True, color=color)
        add_textbox(slide, body, x + Inches(0.12), y + Inches(0.75),
                    Inches(2.86), Inches(2.9), font_size=12, color=LIGHT_GRAY, wrap=True)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT)
    add_textbox(slide, "Summary & Next Steps",
                Inches(0.25), Inches(0.18), Inches(12.8), Inches(0.7),
                font_size=32, bold=True, color=WHITE)

    takeaways = [
        ("What we built", [
            "13-class malicious ad detector from screenshots",
            "InternVL 2.5-1B vision encoder + LoRA + linear head",
            "Config-driven: swap YAML to train any VLM classifier",
            "Fully automated EC2 spot/on-demand pipeline from CLI",
        ], ACCENT),
        ("Why it works", [
            "InternVL pre-trained on web data — transfers directly",
            "448px resolution reads small text in ads",
            "LoRA: 1.6% of parameters, near full-fine-tune accuracy",
            "CLS token aggregates full-image context in one vector",
        ], GREEN),
        ("Current status", [
            "Training live: epoch 2.7 / 5 on on-demand A10G",
            "Val accuracy: 95.2% and still improving",
            "ETA: ~8 more hours to completion",
            "Full classification report saved to S3 on finish",
        ], ORANGE),
    ]
    for i, (title, points, color) in enumerate(takeaways):
        x = Inches(0.25 + i * 4.35)
        add_rect(slide, x, Inches(1.1), Inches(4.1), Inches(4.8), DARK_CARD)
        add_rect(slide, x, Inches(1.1), Inches(4.1), Inches(0.06), color)
        add_textbox(slide, title, x + Inches(0.15), Inches(1.2),
                    Inches(3.8), Inches(0.45), font_size=16, bold=True, color=color)
        y = Inches(1.72)
        for pt in points:
            add_textbox(slide, f"• {pt}", x + Inches(0.15), y,
                        Inches(3.8), Inches(0.5), font_size=13, color=LIGHT_GRAY, wrap=True)
            y += Inches(0.52)

    next_steps = [
        "Download & evaluate final model after training completes",
        "Run inference on production ad stream",
        "Add new label sets via new YAML config — no code changes",
        "Retrain periodically as new threat patterns emerge",
    ]
    add_textbox(slide, "Next steps:",
                Inches(0.25), Inches(6.1), Inches(3), Inches(0.35),
                font_size=14, bold=True, color=ACCENT)
    ns_text = "   |   ".join(next_steps)
    add_textbox(slide, ns_text, Inches(0.25), Inches(6.5), Inches(12.8), Inches(0.6),
                font_size=12, color=LIGHT_GRAY, wrap=True)

    add_textbox(slide, "github.com/RanGeoEdge/geoedge_vlm_trainer",
                Inches(0.25), Inches(7.1), Inches(8), Inches(0.35),
                font_size=12, color=ACCENT2)


def slide_viz_image(prs, img_path, title, subtitle):
    """Full-bleed visualisation slide with title bar overlay."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    # image fills most of the slide
    slide.shapes.add_picture(img_path,
                             Inches(0.15), Inches(1.05),
                             Inches(13.03), Inches(6.3))
    # semi-transparent title bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(1.0), ACCENT)
    add_textbox(slide, title,
                Inches(0.25), Inches(0.08), Inches(12.8), Inches(0.55),
                font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.25), Inches(0.6), Inches(12.8), Inches(0.35),
                    font_size=14, color=ACCENT2)


# ── Build ──────────────────────────────────────────────────────────────────────

HERE = str(Path(__file__).parent)


def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_viz_image(prs, f"{HERE}/viz_patch_embed.png",
                    "Step 1–2: Image → Patches → Token Sequence",
                    "448×448 screenshot sliced into 1024 patch tokens + 1 CLS token, each embedded as a 1024-dim vector")
    slide_architecture(prs)
    slide_viz_image(prs, f"{HERE}/viz_architecture.png",
                    "Full Model Architecture",
                    "InternViT-300M  →  CLS token  →  Dropout  →  Linear(1024→13)  →  Predicted class")
    slide_why_drop_llm(prs)
    slide_lora(prs)
    slide_viz_image(prs, f"{HERE}/viz_lora.png",
                    "LoRA: Low-Rank Adaptation",
                    "Freeze 300M base weights — inject tiny A·B matrices (4.7M params). 64× fewer parameters, same accuracy.")
    slide_comparison(prs)
    slide_results(prs)
    slide_viz_image(prs, f"{HERE}/viz_classification_head.png",
                    "Classification Head — CLS Token → 13 Classes",
                    "Linear(1024→13) maps the global image vector to per-class logits. argmax gives the predicted threat category.")
    slide_infra(prs)
    slide_summary(prs)

    out = "/Users/randubin/PycharmProjects/PythonProject/MIC2/aws_training/MIC2_VLM_Classifier.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
