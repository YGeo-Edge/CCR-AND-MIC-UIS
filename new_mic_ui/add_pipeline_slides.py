#!/usr/bin/env python3
"""
Add two-stage classification pipeline slides:
  1. Pipeline overview (visual flow diagram)
  2. Stage 1 — Binary gate (Benign vs Malicious)
  3. Stage 2 — Malicious type classifier
  4. End-to-end code + decision table
"""
import io, json
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent

PAL = {
    'bg':'#0f172a','card':'#1e293b','border':'#334155',
    'blue':'#38bdf8','white':'#f8fafc','gray':'#94a3b8',
    'green':'#22c55e','yellow':'#f59e0b','red':'#ef4444',
    'orange':'#f97316','teal':'#14b8a6','purple':'#a855f7',
    'dkgreen':'#15803d','dkblue':'#1e3a5f','dkred':'#7f1d1d',
}

CLASS_COL = {
    "fake_av":"#ef4444","financial_scam":"#f97316",
    "misleading_offers":"#eab308","fake_appstore":"#ec4899",
    "tech_support_scam":"#8b5cf6","gift_card_scan":"#f43f5e",
    "forced_notification":"#06b6d4","suspicious_vpn":"#3b82f6",
    "malicious_extension":"#a855f7","fake_updates":"#14b8a6",
    "fake_downloader":"#dc2626",
}

MALICIOUS_CLASSES = list(CLASS_COL.keys())

REC_THRESH = {
    "fake_appstore":0.050,"gift_card_scan":0.350,
    "malicious_extension":0.800,"tech_support_scam":0.990,
    "misleading_offers":0.990,"suspicious_vpn":0.990,"fake_downloader":0.990,
    "fake_av":0.999,"financial_scam":0.999,"forced_notification":0.999,"fake_updates":0.999,
}
REC_RECALL = {
    "fake_appstore":1.00,"gift_card_scan":0.990,
    "malicious_extension":0.705,"tech_support_scam":0.850,
    "misleading_offers":0.890,"suspicious_vpn":0.855,"fake_downloader":0.775,
    "fake_av":0.520,"financial_scam":0.840,"forced_notification":0.410,"fake_updates":0.090,
}

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(PAL['bg'])
    return slide

def txbox(slide, text, x, y, w, h, size=14, bold=False,
          color='white', align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = rgb(color if color.startswith('#') else PAL[color])

def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_id = 5 if radius else 1
    s = slide.shapes.add_shape(shape_id, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line); s.line.width = Pt(1.0)
    else:
        s.line.fill.background()
    return s

def hline(slide, y, x0=0.35, x1=12.98, color='blue', thick=0.035):
    r = slide.shapes.add_shape(1, Inches(x0), Inches(y), Inches(x1-x0), Inches(thick))
    r.fill.solid(); r.fill.fore_color.rgb = rgb(PAL[color]); r.line.fill.background()

def slide_header(slide, title, subtitle=''):
    txbox(slide, title, 0.35, 0.17, 12.6, 0.62, size=26, bold=True)
    if subtitle:
        txbox(slide, subtitle, 0.35, 0.78, 12.6, 0.32, size=12, color='gray')
    hline(slide, 1.09)

def fig_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=150,
                facecolor=PAL['bg'], edgecolor='none')
    buf.seek(0); plt.close(fig); return buf

def add_pic(slide, buf, x, y, w=None, h=None):
    kw = {}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), **kw)

# ═══════════════════════════════════════════════════════════════════════════════
# S1 — Section divider
# ═══════════════════════════════════════════════════════════════════════════════
def slide_section(prs):
    slide = blank_slide(prs)
    rect(slide, 0, 2.75, 13.33, 2.0, '#1e293b')
    txbox(slide, "PART 4", 0.6, 2.15, 12, 0.5, size=14, bold=True, color='blue')
    txbox(slide, "Two-Stage Classification Pipeline", 0.6, 2.68, 12, 0.92, size=36, bold=True)
    txbox(slide, "Stage 1: Binary gate  ·  Stage 2: Malicious type  ·  blank_LP = Benign",
          0.6, 3.72, 12, 0.5, size=16, color='gray')

# ═══════════════════════════════════════════════════════════════════════════════
# S2 — Pipeline overview (big flow diagram)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_pipeline_overview(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Two-Stage Classification Pipeline",
                 "One model — two decisions: first gate malicious vs benign, then identify the threat type")

    fig, ax = plt.subplots(figsize=(13.0, 5.6))
    fig.patch.set_facecolor(PAL['bg'])
    ax.set_facecolor(PAL['bg'])
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.6)
    ax.axis('off')

    def box(cx, cy, w, h, fill, border, label, sublabel='', label_size=11, sub_size=8.5):
        rect_patch = mpatches.FancyBboxPatch(
            (cx - w/2, cy - h/2), w, h,
            boxstyle="round,pad=0.12", linewidth=1.5,
            edgecolor=border, facecolor=fill)
        ax.add_patch(rect_patch)
        ax.text(cx, cy + (0.18 if sublabel else 0), label,
                ha='center', va='center', fontsize=label_size,
                fontweight='bold', color=PAL['white'])
        if sublabel:
            ax.text(cx, cy - 0.30, sublabel, ha='center', va='center',
                    fontsize=sub_size, color=PAL['gray'])

    def arrow(x1, y1, x2, y2, col='#475569', label='', lpos='mid'):
        ax.annotate('', xy=(x2,y2), xytext=(x1,y1),
                    arrowprops=dict(arrowstyle='->', color=col, lw=2.0))
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx+0.05, my+0.12, label, ha='center', fontsize=9,
                    color=col, fontweight='bold')

    # ── INPUT ─────────────────────────────────────────────────────────────────
    box(1.3, 2.8, 1.9, 1.0, '#1e293b', PAL['blue'],
        "Screenshot", "448×448 px")

    arrow(2.25, 2.8, 3.0, 2.8, PAL['blue'])

    # ── MODEL ─────────────────────────────────────────────────────────────────
    box(3.9, 2.8, 1.7, 1.0, '#0f2a4a', PAL['blue'],
        "InternVL\nClassifier", "13 scores out", 10)

    arrow(4.75, 2.8, 5.5, 2.8, PAL['blue'])

    # ── STAGE 1 ───────────────────────────────────────────────────────────────
    box(6.5, 2.8, 1.9, 1.1, '#1a3322', PAL['green'],
        "STAGE 1", "Binary gate", 13)
    ax.text(6.5, 2.38, "Benign+blank_LP vs Malicious", ha='center',
            fontsize=8, color=PAL['gray'])

    # ── BENIGN branch ─────────────────────────────────────────────────────────
    arrow(6.5, 3.35, 6.5, 4.55, PAL['green'], "BENIGN")
    box(6.5, 4.95, 2.6, 0.72, '#0a2a0a', PAL['green'],
        "✅  BENIGN", "Pass — no action needed", 12)
    ax.text(6.5, 4.35, "(incl. blank_LP)", ha='center', fontsize=8, color=PAL['gray'])

    # ── MALICIOUS branch ──────────────────────────────────────────────────────
    arrow(7.45, 2.8, 8.2, 2.8, PAL['red'], "MALICIOUS")

    box(9.2, 2.8, 1.9, 1.1, '#2a0a0a', PAL['red'],
        "STAGE 2", "Malicious type", 13)
    ax.text(9.2, 2.38, "argmax of malicious scores", ha='center',
            fontsize=8, color=PAL['gray'])

    arrow(10.15, 2.8, 10.9, 2.8, PAL['red'])

    # ── OUTPUT box ────────────────────────────────────────────────────────────
    box(11.85, 2.8, 1.9, 2.6, '#1e293b', PAL['orange'],
        "OUTPUT", "", 11)
    threat_types = [
        ("financial_scam",    PAL['orange']),
        ("fake_appstore",     '#ec4899'),
        ("tech_support_scam", '#8b5cf6'),
        ("gift_card_scan",    '#f43f5e'),
        ("fake_av",           PAL['red']),
        ("…8 more",           PAL['gray']),
    ]
    for ti, (name, col) in enumerate(threat_types):
        ax.text(11.85, 3.55 - ti*0.38, name, ha='center', fontsize=7.5,
                color=col, fontweight='bold' if ti < 5 else 'normal')

    # ── confidence threshold note ─────────────────────────────────────────────
    ax.text(6.5, 0.35,
            "Stage 1 threshold: argmax ≠ Benign/blank_LP  →  flag as malicious  (99.4% recall)",
            ha='center', fontsize=9, color=PAL['gray'], style='italic')
    ax.text(9.2, 0.35,
            "Stage 2 threshold: per-class calibrated (FPR < 0.1%)",
            ha='center', fontsize=9, color=PAL['gray'], style='italic')

    fig.tight_layout(pad=0.2)
    add_pic(slide, fig_bytes(fig), 0.15, 1.22, w=13.05)

# ═══════════════════════════════════════════════════════════════════════════════
# S3 — Stage 1 deep-dive
# ═══════════════════════════════════════════════════════════════════════════════
def slide_stage1(prs, binary):
    b = binary['binary']
    slide = blank_slide(prs)
    slide_header(slide, "Stage 1 — Binary Gate: Benign vs Malicious",
                 "Decision: if argmax prediction ∈ {Benign, blank_LP} → pass clean  ·  otherwise → escalate to Stage 2")

    # ── what counts as Benign ─────────────────────────────────────────────────
    rect(slide, 0.30, 1.25, 5.8, 2.55, '#1e293b', PAL['green'])
    txbox(slide, "Benign class =", 0.45, 1.32, 5.5, 0.38, size=13, bold=True, color='green')

    benign_items = [
        ("Benign",   "Legitimate ad / landing page",          PAL['green']),
        ("blank_LP", "Blank or near-empty landing page",      '#64748b'),
    ]
    for i, (cls, desc, col) in enumerate(benign_items):
        ry = 1.72 + i * 0.88
        rect(slide, 0.45, ry, 5.50, 0.76, '#0a2a0a', col)
        txbox(slide, cls,  0.60, ry+0.06, 2.0, 0.35, size=13, bold=True, color=col)
        txbox(slide, desc, 0.60, ry+0.38, 5.0, 0.30, size=10, color='gray')

    txbox(slide, "→ No further action. Image is safe.", 0.45, 3.52, 5.5, 0.32,
          size=11, italic=True, color='green')

    # ── what counts as Malicious ──────────────────────────────────────────────
    rect(slide, 6.35, 1.25, 6.65, 2.55, '#1e293b', PAL['red'])
    txbox(slide, "Malicious class =", 6.50, 1.32, 6.35, 0.38, size=13, bold=True, color='red')

    mal_cols = ['#ef4444','#f97316','#eab308','#ec4899','#8b5cf6','#f43f5e',
                '#06b6d4','#3b82f6','#a855f7','#14b8a6','#dc2626']
    mal_names = ["fake_av","financial_scam","misleading_offers","fake_appstore",
                 "tech_support_scam","gift_card_scan","forced_notification",
                 "suspicious_vpn","malicious_extension","fake_updates","fake_downloader"]
    col_count = 2
    for i, (name, col) in enumerate(zip(mal_names, mal_cols)):
        ci = i % col_count
        ri = i // col_count
        bx = 6.50 + ci * 3.25
        by = 1.72 + ri * 0.38
        rect(slide, bx, by, 3.12, 0.34, '#1e0a0a', col)
        txbox(slide, name, bx+0.10, by+0.04, 3.0, 0.28, size=9, color=col)

    txbox(slide, "→ Escalate to Stage 2 for threat type.", 6.50, 3.52, 6.35, 0.32,
          size=11, italic=True, color='red')

    # ── performance box ───────────────────────────────────────────────────────
    rect(slide, 0.30, 3.95, 12.70, 2.75, '#1e293b', PAL['border'])
    txbox(slide, "Stage 1 performance on test set", 0.45, 4.02, 12.3, 0.38,
          size=13, bold=True, color='blue')

    metrics = [
        ("Accuracy",    f"{b['accuracy']:.2%}",    "overall"),
        ("Recall",      f"{b['recall']:.2%}",       "malicious caught"),
        ("Specificity", f"{b['specificity']:.2%}",  "benign+blank_LP passed"),
        ("Precision",   f"{b['precision']:.2%}",    "flags are real"),
        ("F1",          f"{b['f1']:.2%}",           "harmonic mean"),
    ]
    for i, (lbl, val, sub) in enumerate(metrics):
        cx = 0.45 + i * 2.52
        rect(slide, cx, 4.42, 2.40, 1.35, '#0a2a12', PAL['green'])
        txbox(slide, lbl, cx+0.12, 4.50, 2.18, 0.35, size=10, color='gray')
        txbox(slide, val, cx+0.12, 4.80, 2.18, 0.58, size=24, bold=True, color='green')
        txbox(slide, sub, cx+0.12, 5.28, 2.18, 0.26, size=8.5, color='gray', italic=True)

    # confusion summary
    txbox(slide,
          f"2×2 result:  TP={b['TP']:,}  TN={b['TN']:,}  "
          f"FP={b['FP']} (benign flagged)  FN={b['FN']} (malicious missed)",
          0.45, 5.90, 12.3, 0.35, size=10, color='gray', italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# S4 — Stage 2 deep-dive
# ═══════════════════════════════════════════════════════════════════════════════
def slide_stage2(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Stage 2 — Malicious Type Classifier",
                 "Only runs if Stage 1 flags an image as malicious  ·  identifies which of 11 threat categories")

    # ── class grid with thresholds ────────────────────────────────────────────
    txbox(slide, "11 malicious categories + recommended threshold",
          0.30, 1.25, 8.5, 0.38, size=13, bold=True, color='blue')

    TIERS = [
        ("✅ Deploy — high recall at low FPR", [
            ("fake_appstore",    0.050, 1.00,  "Extremely distinctive visual fingerprint"),
            ("gift_card_scan",   0.350, 0.990, "Tight visual cluster, very reliable"),
            ("tech_support_scam",0.990, 0.850, "Fake browser alert UI"),
            ("misleading_offers",0.990, 0.890, "Deceptive prize / coupon pages"),
            ("financial_scam",   0.999, 0.840, "Needs high threshold to clear FP"),
            ("suspicious_vpn",   0.990, 0.855, "Distinctive VPN push layout"),
        ], PAL['green'], '#0a2a12'),
        ("⚠️  Caution — medium recall", [
            ("fake_downloader",    0.990, 0.775, "Small class, 40 test images"),
            ("malicious_extension",0.800, 0.705, "Confused with misleading_offers"),
            ("fake_av",            0.999, 0.520, "Recall collapses at high threshold"),
        ], PAL['yellow'], '#2a1a00'),
        ("❌  Needs data — low recall", [
            ("forced_notification",0.999, 0.410, "Too visually diverse"),
            ("fake_updates",       0.999, 0.090, "Smallest class (81 test images)"),
        ], PAL['red'], '#1f0808'),
    ]

    y = 1.68
    for tier_label, rows, tcol, tbg in TIERS:
        rect(slide, 0.30, y, 12.70, 0.36, tbg, tcol)
        txbox(slide, tier_label, 0.42, y+0.04, 12.4, 0.28, size=11, bold=True, color=tcol)
        y += 0.36
        cw = [2.4, 1.3, 1.3, 6.5]
        for ri, (cls, t, rec, note) in enumerate(rows):
            rh = 0.37
            bg = '#12201a' if ri%2 else '#0d1a14'
            if tcol == PAL['yellow']:
                bg = '#1a1400' if ri%2 else '#140f00'
            elif tcol == PAL['red']:
                bg = '#1a0808' if ri%2 else '#140606'
            col_c = CLASS_COL.get(cls, '#94a3b8')
            for ci, (v, w) in enumerate(zip(
                [cls, f"t={t:.3f}", f"{rec:.0%}", note], cw
            )):
                rect(slide, 0.30+sum(cw[:ci]), y, w-0.04, rh, bg)
                c = col_c if ci==0 else (tcol if ci==2 else 'white' if ci==1 else 'gray')
                sz = 10 if ci < 3 else 9
                txbox(slide, v, 0.38+sum(cw[:ci]), y+0.05, w-0.14, rh-0.08, size=sz, color=c)
            y += rh
        y += 0.06

    # ── note ──────────────────────────────────────────────────────────────────
    rect(slide, 0.30, y+0.05, 12.70, 0.50, '#1e293b', PAL['border'])
    txbox(slide, "Threshold shown = FPR<0.1% operating point.  "
                 "For review-queue mode use argmax (no threshold) — catches more but sends more to human review.  "
                 "Recall = fraction of true positives caught at that threshold.",
          0.42, y+0.10, 12.45, 0.38, size=9.5, color='gray', italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# S5 — End-to-end code + decision table
# ═══════════════════════════════════════════════════════════════════════════════
def slide_pipeline_code(prs):
    slide = blank_slide(prs)
    slide_header(slide, "End-to-End Pipeline — Code & Decision Table",
                 "One forward pass produces both Stage 1 and Stage 2 decisions simultaneously")

    # ── code ──────────────────────────────────────────────────────────────────
    rect(slide, 0.30, 1.25, 7.20, 5.50, '#0d1117', '#334155')
    txbox(slide, "Python implementation", 0.42, 1.30, 7.0, 0.36, size=11, bold=True, color='blue')

    code = """\
BENIGN_CLASSES = {"Benign", "blank_LP"}

# Per-class thresholds (FPR < 0.1%)
THRESHOLDS = {
    "fake_appstore":        0.050,
    "gift_card_scan":       0.350,
    "malicious_extension":  0.800,
    "tech_support_scam":    0.990,
    "misleading_offers":    0.990,
    "suspicious_vpn":       0.990,
    "fake_downloader":      0.990,
    "fake_av":              0.999,
    "financial_scam":       0.999,
    "forced_notification":  0.999,
    "fake_updates":         0.999,
}

def classify(probs, class_names):
    # ── Stage 1: Binary gate ──────────────────
    pred_cls  = class_names[probs.argmax()]
    pred_conf = probs.max()

    if pred_cls in BENIGN_CLASSES:
        return {
            "stage1": "BENIGN",
            "label":  pred_cls,   # "Benign" or "blank_LP"
            "conf":   pred_conf,
            "stage2": None,
        }

    # ── Stage 2: Malicious type ───────────────
    # Use calibrated per-class threshold
    t = THRESHOLDS.get(pred_cls, 0.500)
    if pred_conf >= t:
        return {
            "stage1": "MALICIOUS",
            "label":  pred_cls,
            "conf":   pred_conf,
            "stage2": "confirmed",
        }
    else:
        # Below threshold → send to human review
        return {
            "stage1": "MALICIOUS",
            "label":  pred_cls,
            "conf":   pred_conf,
            "stage2": "review",
        }"""

    txbox(slide, code, 0.42, 1.65, 7.0, 5.05, size=8.5, color='teal')

    # ── decision table ────────────────────────────────────────────────────────
    txbox(slide, "Decision table", 7.75, 1.25, 5.5, 0.36, size=11, bold=True, color='blue')

    headers = ["Stage 1", "Stage 2", "Action", "Example"]
    cw = [1.5, 1.5, 1.6, 2.9]
    x0, y0, rh = 7.65, 1.65, 0.46

    x = x0
    for col, w in zip(headers, cw):
        rect(slide, x, y0, w-0.04, rh, '#1e3a5f')
        txbox(slide, col, x+0.08, y0+0.08, w-0.14, rh-0.1, size=9, bold=True, color='blue')
        x += w

    rows = [
        ("BENIGN",    "—",          "Pass ✅",        "Benign / blank_LP",       '#0a2a12','green'),
        ("MALICIOUS", "Confirmed",  "Block / Alert ❌","financial_scam ≥ 0.999",  '#2a0a0a','red'),
        ("MALICIOUS", "Review ⚠️",  "Human review",   "financial_scam < 0.999",  '#2a1800','yellow'),
        ("MALICIOUS", "Confirmed",  "Block / Alert ❌","fake_appstore ≥ 0.050",   '#2a0a0a','red'),
        ("MALICIOUS", "Confirmed",  "Block / Alert ❌","fake_updates ≥ 0.999",    '#2a0a0a','red'),
        ("MALICIOUS", "Review ⚠️",  "Human review",   "fake_av < 0.999",         '#2a1800','yellow'),
    ]
    for ri, (s1, s2, action, ex, bg, acol) in enumerate(rows):
        ry = y0 + (ri+1)*rh
        x  = x0
        for ci, (v, w) in enumerate(zip([s1, s2, action, ex], cw)):
            rect(slide, x, ry, w-0.04, rh, bg)
            c = 'green' if s1=='BENIGN' else (acol if ci==2 else 'white')
            txbox(slide, v, x+0.08, ry+0.08, w-0.14, rh-0.1, size=9, color=c)
            x += w

    # ── performance summary ───────────────────────────────────────────────────
    rect(slide, 7.65, 4.60, 5.50, 2.15, '#1e293b', PAL['border'])
    txbox(slide, "Expected pipeline performance", 7.78, 4.67, 5.25, 0.36,
          size=11, bold=True, color='blue')

    perf = [
        ("Stage 1 recall",    "99.4%", "malicious images caught",  'green'),
        ("Stage 1 precision", "99.7%", "flags that are real",       'green'),
        ("Stage 1 FPR",       "0.9%",  "benign images flagged",     'teal'),
        ("Stage 2 accuracy",  "96.1%", "malicious type (13-class)", 'teal'),
    ]
    for pi, (lbl, val, sub, col) in enumerate(perf):
        py = 5.10 + pi * 0.38
        rect(slide, 7.78, py, 5.25, 0.36, '#0d1a25')
        txbox(slide, lbl, 7.88, py+0.05, 2.4, 0.28, size=9, color='gray')
        txbox(slide, val, 10.3, py+0.05, 1.0, 0.28, size=9, bold=True, color=col)
        txbox(slide, sub, 11.3, py+0.05, 1.6, 0.28, size=8.5, color='gray', italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# S6 — Visual example walkthrough
# ═══════════════════════════════════════════════════════════════════════════════
def slide_examples(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Pipeline Walkthrough — Example Outcomes",
                 "Four representative cases showing Stage 1 → Stage 2 flow")

    cases = [
        {
            "title": "Case A — Clean benign page",
            "s1_label": "BENIGN", "s1_conf": "98.2%", "s1_col": PAL['green'],
            "s2_label": "—  No Stage 2 needed", "s2_col": PAL['gray'],
            "action": "✅  Pass — no action",
            "action_col": PAL['green'],
            "scores": {"Benign":0.982,"blank_LP":0.006,"financial_scam":0.004,
                       "fake_av":0.002,"misleading_offers":0.003,"other":0.003},
            "note": "High Benign confidence. Pipeline stops at Stage 1.",
        },
        {
            "title": "Case B — Blank landing page",
            "s1_label": "BENIGN", "s1_conf": "99.5%", "s1_col": PAL['green'],
            "s2_label": "—  blank_LP = Benign", "s2_col": '#64748b',
            "action": "✅  Pass — blank page",
            "action_col": PAL['green'],
            "scores": {"blank_LP":0.995,"Benign":0.003,"financial_scam":0.001,
                       "fake_av":0.0005,"misleading_offers":0.0003,"other":0.0002},
            "note": "blank_LP predicted → treated as Benign. No further action.",
        },
        {
            "title": "Case C — Fake app store (confirmed)",
            "s1_label": "MALICIOUS", "s1_conf": "99.8%", "s1_col": PAL['red'],
            "s2_label": "fake_appstore  (t=0.050, conf≥t)", "s2_col": '#ec4899',
            "action": "❌  Block — fake_appstore",
            "action_col": PAL['red'],
            "scores": {"fake_appstore":0.998,"Benign":0.001,"financial_scam":0.0005,
                       "misleading_offers":0.0003,"fake_av":0.0001,"other":0.0001},
            "note": "Stage 1 flags malicious. Stage 2 confirms fake_appstore at t=0.050 — auto-block.",
        },
        {
            "title": "Case D — Financial scam (review)",
            "s1_label": "MALICIOUS", "s1_conf": "94.1%", "s1_col": PAL['orange'],
            "s2_label": "financial_scam  (t=0.999, conf<t)", "s2_col": PAL['orange'],
            "action": "⚠️  Review — below threshold",
            "action_col": PAL['yellow'],
            "scores": {"financial_scam":0.941,"Benign":0.032,"misleading_offers":0.018,
                       "fake_av":0.005,"fake_updates":0.002,"other":0.002},
            "note": "Stage 1 flags malicious. Stage 2: financial_scam predicted but conf 0.941 < t=0.999 → human review.",
        },
    ]

    for ci, case in enumerate(cases):
        col_off = ci % 2
        row_off = ci // 2
        cx = 0.25 + col_off * 6.62
        cy = 1.28 + row_off * 3.05

        rect(slide, cx, cy, 6.45, 2.90, '#1e293b', case['s1_col'])

        # header
        txbox(slide, case['title'], cx+0.12, cy+0.08, 6.2, 0.35,
              size=11, bold=True, color='white')

        # Stage 1
        rect(slide, cx+0.12, cy+0.48, 2.9, 0.58, '#0d1117', case['s1_col'])
        txbox(slide, "STAGE 1", cx+0.22, cy+0.50, 2.7, 0.26, size=8, color='gray', bold=True)
        txbox(slide, f"{case['s1_label']}  {case['s1_conf']}",
              cx+0.22, cy+0.72, 2.7, 0.28, size=11, bold=True, color=case['s1_col'])

        # Stage 2
        rect(slide, cx+3.20, cy+0.48, 3.1, 0.58, '#0d1117', case['s2_col'])
        txbox(slide, "STAGE 2", cx+3.30, cy+0.50, 2.9, 0.26, size=8, color='gray', bold=True)
        txbox(slide, case['s2_label'],
              cx+3.30, cy+0.72, 2.9, 0.28, size=10, bold=True, color=case['s2_col'])

        # action
        rect(slide, cx+0.12, cy+1.14, 6.20, 0.42, '#0d1117', case['action_col'])
        txbox(slide, case['action'],
              cx+0.22, cy+1.18, 6.0, 0.32, size=12, bold=True,
              color=case['action_col'], align=PP_ALIGN.CENTER)

        # note
        txbox(slide, case['note'],
              cx+0.12, cy+1.62, 6.20, 0.55, size=9, color='gray', italic=True)

        # mini score bars
        fig, ax = plt.subplots(figsize=(2.8, 1.55))
        fig.patch.set_facecolor('#1e293b'); ax.set_facecolor('#0d1117')
        labels = list(case['scores'].keys())
        vals   = list(case['scores'].values())
        cols   = [PAL['green'] if l in ('Benign','blank_LP')
                  else CLASS_COL.get(l, '#475569') for l in labels]
        ax.barh(labels, vals, color=cols, height=0.6)
        ax.set_xlim(0, 1.05)
        ax.tick_params(colors=PAL['gray'], labelsize=6)
        ax.spines[:].set_color('#334155')
        ax.grid(axis='x', color='#334155', lw=0.3)
        fig.tight_layout(pad=0.3)
        add_pic(slide, fig_bytes(fig), cx+0.12, cy+2.22, w=6.20)

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    bin_path = HERE / "binary_validation_results.json"
    if not bin_path.exists():
        print("ERROR: binary_validation_results.json not found.")
        raise SystemExit(1)

    with open(bin_path) as f:
        binary = json.load(f)

    pptx_path = HERE / "MIC2_VLM_Classifier.pptx"
    prs = Presentation(str(pptx_path))
    print(f"Loaded presentation: {len(prs.slides)} existing slides")

    # Insert after slide 14 (original content) — before Part 2
    # We'll append to end — order: 14 original + 11 eval + 7 benchmark + 56 errors
    # New pipeline slides go right after slide 14 (index 14)
    # Easiest: build a temp prs with just these slides, then splice
    # For simplicity, append at position 15 by building them first in a temp file
    # then inserting — python-pptx doesn't support insert, so we rebuild

    print("Building pipeline slides in correct position …")
    from pptx import Presentation as Prs
    from copy import deepcopy
    from lxml import etree

    # Build pipeline slides in a temp presentation
    tmp_prs = Prs()
    tmp_prs.slide_width  = prs.slide_width
    tmp_prs.slide_height = prs.slide_height

    # We need a blank layout in tmp_prs — copy from main prs layout
    # Simpler: just append to the end of the current presentation
    slide_section(prs);          print("  1/5 section divider")
    slide_pipeline_overview(prs);print("  2/5 pipeline overview")
    slide_stage1(prs, binary);   print("  3/5 Stage 1")
    slide_stage2(prs);           print("  4/5 Stage 2")
    slide_pipeline_code(prs);    print("  5/5 code + decision table")
    slide_examples(prs);         print("  6/6 example walkthrough")

    prs.save(str(pptx_path))
    print(f"\nSaved → {pptx_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
